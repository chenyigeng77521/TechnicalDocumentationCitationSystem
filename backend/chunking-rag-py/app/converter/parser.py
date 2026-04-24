from pathlib import Path
from typing import Any

import fitz  # PyMuPDF
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def parse(path: Path) -> tuple[str, dict[int, Any]]:
    """按扩展名分派。返回 (markdown_text, line_map)。line_map MVP 返回 {}。"""
    ext = path.suffix.lower()
    if ext == ".md":
        return _parse_md(path)
    if ext == ".pdf":
        return _parse_pdf(path)
    if ext == ".docx":
        return _parse_docx(path)
    if ext == ".pptx":
        return _parse_pptx(path)
    if ext == ".xlsx":
        return _parse_xlsx(path)
    raise ValueError(f"unsupported file format: {ext}")


def _parse_md(path: Path) -> tuple[str, dict[int, Any]]:
    return path.read_text(encoding="utf-8"), {}


def _parse_pdf(path: Path) -> tuple[str, dict[int, Any]]:
    doc = fitz.open(path)
    parts: list[str] = []
    try:
        for i, page in enumerate(doc, 1):
            text = page.get_text()
            if text.strip():
                parts.append(f"## Page {i}\n\n{text.strip()}")
    finally:
        doc.close()
    return "\n\n".join(parts), {}


def _parse_docx(path: Path) -> tuple[str, dict[int, Any]]:
    doc = Document(path)
    lines: list[str] = []
    for p in doc.paragraphs:
        style = (p.style.name or "").lower() if p.style else ""
        text = p.text.strip()
        if not text:
            continue
        if style.startswith("heading 1"):
            lines.append(f"# {text}")
        elif style.startswith("heading 2"):
            lines.append(f"## {text}")
        elif style.startswith("heading"):
            lines.append(f"### {text}")
        else:
            lines.append(text)
    return "\n\n".join(lines), {}


def _parse_pptx(path: Path) -> tuple[str, dict[int, Any]]:
    pres = Presentation(path)
    parts: list[str] = []
    for i, slide in enumerate(pres.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    parts.append(text)
    return "\n\n".join(parts), {}


def _parse_xlsx(path: Path) -> tuple[str, dict[int, Any]]:
    wb = load_workbook(path, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"## Sheet {ws.title}")
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        header = rows[0]
        parts.append("| " + " | ".join(str(c or "") for c in header) + " |")
        parts.append("|" + "|".join(["---"] * len(header)) + "|")
        for row in rows[1:]:
            parts.append("| " + " | ".join(str(c or "") for c in row) + " |")
    return "\n\n".join(parts), {}
