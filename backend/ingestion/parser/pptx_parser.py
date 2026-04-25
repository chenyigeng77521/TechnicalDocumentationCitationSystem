"""PPTX 解析器。每 slide → 一个段落。"""
from pathlib import Path
from pptx import Presentation
from backend.ingestion.parser.types import ParseResult


async def parse(path: Path) -> ParseResult:
    p = Presentation(path)
    slide_texts = []
    for i, slide in enumerate(p.slides):
        chunks = [f"### Slide {i + 1}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        chunks.append(text)
        slide_texts.append("\n".join(chunks))
    return ParseResult(
        raw_text="\n\n".join(slide_texts),
        title_tree=[],
        content_type="document",
        metadata={"slide_count": len(p.slides)},
    )
