"""PPTX 解析器。每 slide → 一个段落，slide 名进 title_tree。"""
from pathlib import Path
from pptx import Presentation
from backend.ingestion.parser.types import ParseResult, TitleNode


async def parse(path: Path) -> ParseResult:
    p = Presentation(path)
    slide_texts: list[str] = []
    slide_titles: list[str] = []  # 用作 title_tree 的文本

    for i, slide in enumerate(p.slides):
        slide_no = i + 1
        # 优先用 slide 第一个文本框作为标题，没有就用 "Slide N"
        slide_label = f"Slide {slide_no}"
        first_text_in_slide = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        if first_text_in_slide is None:
                            first_text_in_slide = text
        if first_text_in_slide:
            # 用 slide 的第一段文字作为更有意义的 title（截断 50 字防过长）
            slide_label = first_text_in_slide[:50]

        slide_titles.append(slide_label)

        chunks = [f"### Slide {slide_no}: {slide_label}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        chunks.append(text)
        slide_texts.append("\n".join(chunks))

    raw_text = "\n\n".join(slide_texts)

    # 给每个 slide 标题加 TitleNode（level=3，对应 "### Slide N: xxx"）
    title_nodes: list[TitleNode] = []
    cursor = 0
    for i, slide_label in enumerate(slide_titles):
        header = f"### Slide {i + 1}: {slide_label}"
        offset = raw_text.find(header, cursor)
        if offset >= 0:
            title_nodes.append(TitleNode(
                level=3,
                text=f"Slide {i + 1}: {slide_label}",
                char_offset=offset,
            ))
            cursor = offset + len(header)

    return ParseResult(
        raw_text=raw_text,
        title_tree=title_nodes,
        content_type="document",
        metadata={"slide_count": len(p.slides)},
    )
