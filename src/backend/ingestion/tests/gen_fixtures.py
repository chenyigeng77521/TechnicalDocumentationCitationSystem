"""Generate fixture files for parser tests. Run once with: python gen_fixtures.py"""
import pathlib

FIXTURES = pathlib.Path(__file__).parent / "fixtures"
FIXTURES.mkdir(exist_ok=True)


def gen_pdf():
    import fitz
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 100), "Hello PDF World", fontsize=14)
    # Use built-in CJK font
    page.insert_text((50, 130), "这是一段中文。", fontsize=14, fontname="china-s")
    out = FIXTURES / "sample.pdf"
    doc.save(str(out))
    doc.close()
    print(f"Generated {out}")

    # verify
    doc2 = fitz.open(str(out))
    text = doc2[0].get_text()
    doc2.close()
    print(f"Extracted text: {repr(text)}")


def gen_docx():
    from docx import Document
    d = Document()
    d.add_heading("DocX Title", level=1)
    d.add_paragraph("First paragraph")
    d.add_heading("Sub", level=2)
    d.add_paragraph("Sub content")
    out = FIXTURES / "sample.docx"
    d.save(str(out))
    print(f"Generated {out}")


def gen_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Header"
    ws["A2"] = "Value"
    out = FIXTURES / "sample.xlsx"
    wb.save(str(out))
    print(f"Generated {out}")


def gen_pptx():
    from pptx import Presentation
    p = Presentation()
    slide = p.slides.add_slide(p.slide_layouts[5])
    slide.shapes.title.text = "Slide One"
    out = FIXTURES / "sample.pptx"
    p.save(str(out))
    print(f"Generated {out}")


if __name__ == "__main__":
    gen_pdf()
    gen_docx()
    gen_xlsx()
    gen_pptx()
    print("All fixtures generated.")
